"""
Week 1 - Introduction to Generative Deep Learning
Slide Generator for SEAS 8525: Computer Vision and Generative AI
Textbook: Generative Deep Learning by David Foster (O'Reilly, 2nd ed.)
Light academic theme - white backgrounds, dark navy headers, blue accents
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Light Academic Color Palette ─────────────────────────────────────────────
SLIDE_BG  = RGBColor(0xFF, 0xFF, 0xFF)  # white slide background
CARD_BG   = RGBColor(0xF0, 0xF5, 0xFF)  # very light blue - card background
HDR_BG    = RGBColor(0x1A, 0x3A, 0x5C)  # dark navy - header
BLUE      = RGBColor(0x00, 0x77, 0xB6)  # primary blue accent
BLUE_LT   = RGBColor(0x90, 0xC8, 0xE8)  # light blue
TEAL      = RGBColor(0x00, 0x83, 0x8F)  # dark teal
TEXT      = RGBColor(0x1A, 0x2A, 0x3A)  # main body text (near-black)
TEXT_S    = RGBColor(0x54, 0x6E, 0x7A)  # secondary text
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BORDER    = RGBColor(0xB0, 0xC4, 0xD8)  # light blue-gray border
GREEN_D   = RGBColor(0x00, 0x69, 0x5C)  # dark green
GREEN_L   = RGBColor(0xE0, 0xF2, 0xF1)  # light green bg
RED_D     = RGBColor(0xC6, 0x28, 0x28)  # dark red
RED_L     = RGBColor(0xFF, 0xEB, 0xEE)  # light red bg
ORG_D     = RGBColor(0xBF, 0x36, 0x0C)  # dark orange
ORG_L     = RGBColor(0xFF, 0xF3, 0xE0)  # light orange bg
PUR_D     = RGBColor(0x4A, 0x27, 0xA0)  # dark purple
PUR_L     = RGBColor(0xF3, 0xE5, 0xF5)  # light purple bg
AMB       = RGBColor(0xFF, 0x6F, 0x00)  # amber/orange
AMB_L     = RGBColor(0xFF, 0xF8, 0xE1)  # light amber bg
PINK_D    = RGBColor(0xAD, 0x14, 0x57)  # dark rose/pink
PINK_L    = RGBColor(0xFC, 0xE4, 0xEC)  # light pink bg
NAVY_LT   = RGBColor(0xE8, 0xF0, 0xF8)  # very light navy

# ── Slide dimensions (widescreen 16:9) ───────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)


# ── Helper utilities ──────────────────────────────────────────────────────────

def new_slide(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = SLIDE_BG
    return slide


def add_rect(slide, x, y, w, h, fill=None, line=None, lw=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = lw
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h, text, size=Pt(16), bold=False, color=None,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    if color is None:
        color = TEXT
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb


def add_multiline(slide, x, y, w, h, lines, size=Pt(15), color=None):
    if color is None:
        color = TEXT
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        if isinstance(line, dict):
            run.text = line.get("text", "")
            run.font.size = line.get("size", size)
            run.font.bold = line.get("bold", False)
            run.font.color.rgb = line.get("color", color)
            run.font.italic = line.get("italic", False)
        else:
            run.text = line
            run.font.size = size
            run.font.color.rgb = color
        run.font.name = "Calibri"
    return txb


def header_bar(slide, title, subtitle=None):
    """Dark navy header with blue accent strip."""
    add_rect(slide, 0, 0, W, Inches(0.06), fill=BLUE)
    add_rect(slide, 0, Inches(0.06), W, Inches(1.05), fill=HDR_BG)
    add_textbox(slide, Inches(0.4), Inches(0.1), Inches(12.0), Inches(0.7),
                title, size=Pt(30), bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.4), Inches(0.78), Inches(12.0), Inches(0.35),
                    subtitle, size=Pt(15), color=BLUE_LT, italic=True)


def footer(slide, text="SEAS 8525 - Generative Deep Learning | Week 1"):
    add_rect(slide, 0, H - Inches(0.38), W, Inches(0.38), fill=NAVY_LT)
    add_rect(slide, 0, H - Inches(0.38), W, Inches(0.025), fill=BLUE)
    add_textbox(slide, Inches(0.4), H - Inches(0.36), Inches(10), Inches(0.3),
                text, size=Pt(11), color=TEXT_S)


def info_card(slide, x, y, w, h, title, body_lines, accent, light_bg,
              body_size=Pt(13.5)):
    """A card with light background and colored left/top accent."""
    add_rect(slide, x, y, w, h, fill=light_bg, line=accent, lw=Pt(1.2))
    add_rect(slide, x, y, Inches(0.1), h, fill=accent)
    add_textbox(slide, x + Inches(0.18), y + Inches(0.1), w - Inches(0.28), Inches(0.45),
                title, size=Pt(16), bold=True, color=accent)
    add_multiline(slide, x + Inches(0.18), y + Inches(0.6), w - Inches(0.28),
                  h - Inches(0.75), body_lines, size=body_size, color=TEXT)


def eq_box(slide, x, y, w, h, text, color=BLUE):
    """Equation box with light blue background."""
    add_rect(slide, x, y, w, h, fill=CARD_BG, line=color, lw=Pt(1.5))
    add_rect(slide, x, y, Inches(0.1), h, fill=color)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.1),
                w - Inches(0.3), h - Inches(0.2),
                text, size=Pt(17), bold=True, color=color, italic=True,
                align=PP_ALIGN.CENTER)


def divider(slide, y):
    add_rect(slide, Inches(0.4), y, W - Inches(0.8), Inches(0.025), fill=BLUE)


# ── Slide 1: Title ────────────────────────────────────────────────────────────
def slide_title(prs):
    s = new_slide(prs)
    # Left blue accent bar
    add_rect(s, 0, 0, Inches(0.18), H, fill=HDR_BG)
    add_rect(s, 0, 0, Inches(0.08), H, fill=BLUE)
    # Course label
    add_rect(s, Inches(0.4), Inches(0.5), Inches(6.5), Inches(0.55), fill=HDR_BG)
    add_textbox(s, Inches(0.5), Inches(0.55), Inches(6.3), Inches(0.45),
                "SEAS 8525: Computer Vision and Generative AI",
                size=Pt(16), color=WHITE, bold=True)
    # Main title
    add_textbox(s, Inches(0.4), Inches(1.4), Inches(12.5), Inches(0.7),
                "Introduction to", size=Pt(36), color=TEXT_S)
    add_textbox(s, Inches(0.4), Inches(2.0), Inches(12.5), Inches(1.5),
                "Generative Deep\nLearning",
                size=Pt(60), bold=True, color=HDR_BG)
    # Chapter tag
    add_rect(s, Inches(0.4), Inches(4.0), Inches(5.0), Inches(0.6), fill=BLUE)
    add_textbox(s, Inches(0.5), Inches(4.05), Inches(4.8), Inches(0.5),
                "Week 1  |  Chapter 1: What Is Generative Modeling?",
                size=Pt(15), bold=True, color=WHITE)
    add_rect(s, Inches(0.4), Inches(4.75), Inches(7), Inches(0.025), fill=BLUE)
    add_textbox(s, Inches(0.4), Inches(4.9), Inches(9), Inches(0.4),
                "Textbook: Generative Deep Learning - David Foster (O'Reilly, 2nd ed.)",
                size=Pt(14), color=TEXT_S)
    add_textbox(s, Inches(0.4), Inches(5.4), Inches(9), Inches(0.4),
                "Spring 2026", size=Pt(14), color=TEXT_S)
    footer(s)
    return s


# ── Slide 2: Learning Objectives ─────────────────────────────────────────────
def slide_objectives(prs):
    s = new_slide(prs)
    header_bar(s, "Learning Objectives",
               "By the end of this lecture you will be able to:")
    footer(s)
    objectives = [
        (BLUE,   "01", "Define Generative Modeling",
         "Articulate the formal difference between discriminative and generative models "
         "using probability notation p(y|x) vs p(x)."),
        (GREEN_D, "02", "Apply Core Probability Theory",
         "Use PDFs, parametric families, and Maximum Likelihood Estimation (MLE) "
         "to describe how models are fitted to data."),
        (ORG_D,  "03", "Explain Latent Variables",
         "Describe how latent variables z encode hidden structure, "
         "and derive the marginal p(x) = integral of p(x|z) p(z) dz."),
        (PUR_D,  "04", "Compare Generative Model Families",
         "Identify and contrast the 6 major families: "
         "EBMs, VAEs, GANs, Flows, Diffusion, and Autoregressive models."),
    ]
    xs = [Inches(0.3), Inches(6.85)]
    ys = [Inches(1.4), Inches(4.2)]
    for i, (color, num, title, body) in enumerate(objectives):
        cx = xs[i % 2]
        cy = ys[i // 2]
        # card
        light = GREEN_L if color == GREEN_D else (ORG_L if color == ORG_D else (PUR_L if color == PUR_D else CARD_BG))
        add_rect(s, cx, cy, Inches(6.2), Inches(2.55), fill=light, line=color, lw=Pt(1.5))
        add_rect(s, cx, cy, Inches(0.65), Inches(2.55), fill=color)
        add_textbox(s, cx, cy + Inches(0.8), Inches(0.65), Inches(0.6),
                    num, size=Pt(20), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, cx + Inches(0.75), cy + Inches(0.1), Inches(5.3), Inches(0.5),
                    title, size=Pt(17), bold=True, color=color)
        add_textbox(s, cx + Inches(0.75), cy + Inches(0.7), Inches(5.3), Inches(1.65),
                    body, size=Pt(13.5), color=TEXT, wrap=True)
    return s


# ── Slide 3: The AI Landscape ─────────────────────────────────────────────────
def slide_ai_landscape(prs):
    s = new_slide(prs)
    header_bar(s, "Where Does Generative AI Fit?", "The AI Landscape")
    footer(s)
    boxes = [
        (Inches(0.35), Inches(1.4), Inches(12.6), Inches(5.5),
         RGBColor(0xF0, 0xF4, 0xF8), "Artificial Intelligence",    Pt(16), TEXT_S),
        (Inches(0.7),  Inches(1.7), Inches(11.6), Inches(4.8),
         RGBColor(0xE3, 0xEE, 0xF8), "Machine Learning",           Pt(16), TEXT_S),
        (Inches(1.1),  Inches(2.05), Inches(10.6), Inches(4.0),
         RGBColor(0xD0, 0xE8, 0xF5), "Deep Learning",              Pt(16), HDR_BG),
        (Inches(2.0),  Inches(2.6), Inches(8.5),  Inches(2.7),
         RGBColor(0xB8, 0xD8, 0xF0), "Generative Deep Learning",   Pt(18), HDR_BG),
    ]
    for bx, by, bw, bh, bc, label, fs, fc in boxes:
        add_rect(s, bx, by, bw, bh, fill=bc, line=BLUE, lw=Pt(0.8))
        add_textbox(s, bx + Inches(0.15), by + Inches(0.1), bw - Inches(0.3), Inches(0.45),
                    label, size=fs, bold=True, color=fc)
    examples = ["VAEs", "GANs", "Diffusion", "Flows", "Autoregressive"]
    colors = [TEAL, GREEN_D, PINK_D, ORG_D, PUR_D]
    ex_x = Inches(2.3)
    for i, (ex, ec) in enumerate(zip(examples, colors)):
        add_rect(s, ex_x + Inches(i * 1.55), Inches(3.45), Inches(1.4), Inches(0.6), fill=ec)
        add_textbox(s, ex_x + Inches(i * 1.55) + Inches(0.05), Inches(3.5),
                    Inches(1.3), Inches(0.45), ex, size=Pt(13), bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.35),
                "Generative models learn the distribution of the data itself - they create new, realistic samples.",
                size=Pt(13), color=AMB, italic=True)
    return s


# ── Slide 4: Discriminative vs Generative ─────────────────────────────────────
def slide_disc_vs_gen(prs):
    s = new_slide(prs)
    header_bar(s, "Discriminative vs. Generative Models")
    footer(s)
    # Left - Discriminative
    add_rect(s, Inches(0.3), Inches(1.35), Inches(5.9), Inches(5.6),
             fill=CARD_BG, line=BLUE, lw=Pt(1.2))
    add_rect(s, Inches(0.3), Inches(1.35), Inches(5.9), Inches(0.6), fill=HDR_BG)
    add_textbox(s, Inches(0.5), Inches(1.4), Inches(5.5), Inches(0.5),
                "Discriminative Model", size=Pt(20), bold=True, color=WHITE)
    disc_lines = [
        {"text": "Models:  p(y | x)", "size": Pt(19), "bold": True, "color": BLUE},
        {"text": "", "size": Pt(6), "color": TEXT},
        {"text": "Given input x, predict output label y", "size": Pt(14), "color": TEXT},
        {"text": "Learns a decision boundary in input space", "size": Pt(14), "color": TEXT},
        {"text": "Cannot generate new data samples", "size": Pt(14), "color": TEXT_S},
        {"text": "", "size": Pt(8), "color": TEXT},
        {"text": "Examples:", "size": Pt(14), "bold": True, "color": HDR_BG},
        {"text": "Image classifiers, object detectors,", "size": Pt(13), "color": TEXT_S},
        {"text": "sentiment analysis, fraud detection", "size": Pt(13), "color": TEXT_S},
        {"text": "", "size": Pt(8), "color": TEXT},
        {"text": "Analogy: A critic who judges art", "size": Pt(13), "italic": True, "color": AMB},
    ]
    add_multiline(s, Inches(0.5), Inches(2.1), Inches(5.5), Inches(4.5), disc_lines)
    # VS badge
    add_rect(s, Inches(6.1), Inches(3.6), Inches(0.8), Inches(0.8), fill=PINK_D)
    add_textbox(s, Inches(6.1), Inches(3.62), Inches(0.8), Inches(0.65),
                "VS", size=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Right - Generative
    add_rect(s, Inches(7.0), Inches(1.35), Inches(6.0), Inches(5.6),
             fill=GREEN_L, line=GREEN_D, lw=Pt(1.2))
    add_rect(s, Inches(7.0), Inches(1.35), Inches(6.0), Inches(0.6), fill=GREEN_D)
    add_textbox(s, Inches(7.2), Inches(1.4), Inches(5.6), Inches(0.5),
                "Generative Model", size=Pt(20), bold=True, color=WHITE)
    gen_lines = [
        {"text": "Models:  p(x)  or  p(x | y)", "size": Pt(19), "bold": True, "color": GREEN_D},
        {"text": "", "size": Pt(6), "color": TEXT},
        {"text": "Models the full distribution of data x", "size": Pt(14), "color": TEXT},
        {"text": "Can generate brand-new, unseen samples", "size": Pt(14), "color": TEXT},
        {"text": "Understands the underlying structure of data", "size": Pt(14), "color": TEXT},
        {"text": "", "size": Pt(8), "color": TEXT},
        {"text": "Examples:", "size": Pt(14), "bold": True, "color": HDR_BG},
        {"text": "DALL-E, GPT, Stable Diffusion,", "size": Pt(13), "color": TEXT_S},
        {"text": "VAEs, GANs, Diffusion Models", "size": Pt(13), "color": TEXT_S},
        {"text": "", "size": Pt(8), "color": TEXT},
        {"text": "Analogy: An artist who creates art", "size": Pt(13), "italic": True, "color": AMB},
    ]
    add_multiline(s, Inches(7.2), Inches(2.1), Inches(5.6), Inches(4.5), gen_lines)
    return s


# ── Slide 5: Formal Problem Statement ─────────────────────────────────────────
def slide_formal_problem(prs):
    s = new_slide(prs)
    header_bar(s, "The Generative Modeling Problem", "Formal Definition")
    footer(s)
    eq_box(s, Inches(1.0), Inches(1.45), Inches(11.3), Inches(0.9),
           "Given: observations X = {x_1, x_2, ..., x_n} drawn from unknown distribution p_data(x)",
           color=HDR_BG)
    eq_box(s, Inches(1.0), Inches(2.5), Inches(11.3), Inches(0.8),
           "Goal: build p_model(x ; theta) that approximates p_data(x)",
           color=BLUE)
    add_textbox(s, Inches(0.4), Inches(3.55), Inches(12.5), Inches(0.38),
                "Two requirements for a valid generative model:", size=Pt(16), color=HDR_BG, bold=True)
    divider(s, Inches(3.95))
    reqs = [
        (Inches(0.3), "Requirement 1: Accuracy",
         "p_model must assign HIGH probability to data that looks real,\n"
         "and LOW probability to implausible observations.",
         BLUE, CARD_BG),
        (Inches(7.0), "Requirement 2: Generativity",
         "We must be able to SAMPLE new observations  x' ~ p_model(x)\n"
         "that are realistic and diverse, not just memorized training data.",
         GREEN_D, GREEN_L),
    ]
    for rx, rtitle, rbody, rc, rlight in reqs:
        add_rect(s, rx, Inches(4.1), Inches(6.25), Inches(3.0), fill=rlight, line=rc, lw=Pt(1.5))
        add_rect(s, rx, Inches(4.1), Inches(0.1), Inches(3.0), fill=rc)
        add_textbox(s, rx + Inches(0.2), Inches(4.2), Inches(5.9), Inches(0.5),
                    rtitle, size=Pt(17), bold=True, color=rc)
        add_textbox(s, rx + Inches(0.2), Inches(4.8), Inches(5.9), Inches(2.1),
                    rbody, size=Pt(14.5), color=TEXT, wrap=True)
    return s


# ── Slide 6: Applications ─────────────────────────────────────────────────────
def slide_applications(prs):
    s = new_slide(prs)
    header_bar(s, "Why Does This Matter?",
               "Real-World Applications of Generative Deep Learning")
    footer(s)
    apps = [
        ("Image Generation",     "DALL-E 3, Midjourney, Stable Diffusion\nPhotorealistic images from text prompts",     BLUE,    Inches(0.25)),
        ("Text and Language",    "GPT-4, Claude, Llama\nHuman-quality text generation and reasoning",                   GREEN_D, Inches(3.5)),
        ("Music and Audio",      "MusicLM, AudioCraft, Suno\nCompose original music and synthesize speech",              ORG_D,   Inches(6.75)),
        ("Video Generation",     "Sora, Gen-3, Runway\nFull HD video from text descriptions",                           PINK_D,  Inches(10.0)),
        ("Drug Discovery",       "AlphaFold, molecule generation\nDesign novel proteins and drug compounds",              TEAL,    Inches(0.25)),
        ("Scientific Simulation","Climate, physics, genomics\nGenerate plausible scientific scenarios",                  PUR_D,   Inches(3.5)),
        ("Art and Creativity",   "Style transfer, inpainting, upscaling\nAI as a creative collaborator",                GREEN_D, Inches(6.75)),
        ("3D and World Models",  "NeRF, Gaussian Splatting\nGenerate 3D scenes for VR/AR applications",                 ORG_D,   Inches(10.0)),
    ]
    lights = [CARD_BG, GREEN_L, ORG_L, PINK_L, CARD_BG, PUR_L, GREEN_L, ORG_L]
    for i, (title, body, color, x) in enumerate(apps):
        y = Inches(1.5) if i < 4 else Inches(4.35)
        w = Inches(3.1)
        h = Inches(2.65)
        light = lights[i]
        add_rect(s, x, y, w, h, fill=light, line=color, lw=Pt(1.5))
        add_rect(s, x, y, w, Inches(0.12), fill=color)
        add_textbox(s, x + Inches(0.15), y + Inches(0.18), w - Inches(0.3), Inches(0.45),
                    title, size=Pt(15), bold=True, color=color)
        add_textbox(s, x + Inches(0.15), y + Inches(0.72), w - Inches(0.3), Inches(1.75),
                    body, size=Pt(12.5), color=TEXT, wrap=True)
    return s


# ── Slide 7: First Generative Model - Setup ───────────────────────────────────
def slide_first_model_setup(prs):
    s = new_slide(prs)
    header_bar(s, "Our First Generative Model", "Step 1: The Setup")
    footer(s)
    add_textbox(s, Inches(0.4), Inches(1.35), Inches(12.5), Inches(0.45),
                "Suppose we have a small dataset of observations - for example, a collection of artworks:",
                size=Pt(15), color=TEXT)
    colors = [BLUE, GREEN_D, PINK_D, ORG_D, TEAL]
    lights = [CARD_BG, GREEN_L, PINK_L, ORG_L, CARD_BG]
    labels = ["x_1", "x_2", "x_3", "x_4", "x_5"]
    for i in range(5):
        bx = Inches(0.5 + i * 2.4)
        add_rect(s, bx, Inches(2.0), Inches(2.1), Inches(2.1), fill=lights[i], line=colors[i], lw=Pt(1.5))
        add_rect(s, bx + Inches(0.15), Inches(2.2), Inches(1.8), Inches(1.1),
                 fill=RGBColor(0xB0 + i * 8, 0xC8 + i * 4, 0xE0))
        add_textbox(s, bx, Inches(3.35), Inches(2.1), Inches(0.4),
                    labels[i], size=Pt(16), bold=True, color=colors[i],
                    align=PP_ALIGN.CENTER, italic=True)
    add_textbox(s, Inches(12.2), Inches(2.8), Inches(0.8), Inches(0.5),
                "...", size=Pt(28), bold=True, color=TEXT_S)
    add_rect(s, Inches(0.4), Inches(4.35), Inches(12.5), Inches(1.55),
             fill=CARD_BG, line=BLUE, lw=Pt(1.5))
    add_rect(s, Inches(0.4), Inches(4.35), Inches(0.1), Inches(1.55), fill=BLUE)
    add_textbox(s, Inches(0.65), Inches(4.45), Inches(12.0), Inches(0.5),
                "Key Question: What rules govern what a valid observation looks like?",
                size=Pt(17), bold=True, color=HDR_BG)
    add_textbox(s, Inches(0.65), Inches(5.0), Inches(12.0), Inches(0.6),
                "There exists some unknown distribution p_data(x) that generated all these observations. "
                "Our goal: learn a parametric model p_model(x; theta) that approximates this distribution.",
                size=Pt(14), color=TEXT, wrap=True)
    return s


# ── Slide 8: First Generative Model - Learning ────────────────────────────────
def slide_first_model_learning(prs):
    s = new_slide(prs)
    header_bar(s, "Our First Generative Model", "Step 2: Learning the Distribution")
    footer(s)
    steps = [
        ("1. Choose a Model Family",
         "Select a parametric family of distributions:\ne.g., Gaussian, mixture of Gaussians, neural network.\nThis defines the space of possible models.",
         BLUE, CARD_BG, Inches(0.3), Inches(1.5)),
        ("2. Estimate Parameters theta",
         "Find parameters theta* that maximize the probability\nof the observed data:\ntheta* = argmax  sum_i  log p(x_i ; theta)",
         GREEN_D, GREEN_L, Inches(0.3), Inches(3.85)),
        ("3. Evaluate the Model",
         "Check that p_model assigns high probability to real\nobservations and low probability to implausible ones.\nUse held-out data to test generalization.",
         ORG_D, ORG_L, Inches(6.85), Inches(1.5)),
        ("4. Refine and Iterate",
         "If the model does not fit well, try a more expressive\nfamily. Deep neural networks give us enormous flexibility\nfor complex, high-dimensional data.",
         PUR_D, PUR_L, Inches(6.85), Inches(3.85)),
    ]
    for title, body, color, light, x, y in steps:
        add_rect(s, x, y, Inches(6.25), Inches(2.2), fill=light, line=color, lw=Pt(1.5))
        add_rect(s, x, y, Inches(0.1), Inches(2.2), fill=color)
        add_textbox(s, x + Inches(0.2), y + Inches(0.1), Inches(5.9), Inches(0.48),
                    title, size=Pt(16), bold=True, color=color)
        add_textbox(s, x + Inches(0.2), y + Inches(0.65), Inches(5.9), Inches(1.4),
                    body, size=Pt(13.5), color=TEXT, wrap=True)
    return s


# ── Slide 9: First Generative Model - Sampling ───────────────────────────────
def slide_first_model_sampling(prs):
    s = new_slide(prs)
    header_bar(s, "Our First Generative Model", "Step 3: Generating New Observations")
    footer(s)
    add_textbox(s, Inches(0.4), Inches(1.4), Inches(12.5), Inches(0.4),
                "Once we have a fitted model p_model(x ; theta*), we can generate brand-new observations:",
                size=Pt(15), color=TEXT)
    flow = [
        (Inches(0.4),  "Sample  z ~ p(z)", BLUE,   "Latent code / noise"),
        (Inches(4.0),  "Generator  G(z ; theta*)", GREEN_D, "Learned neural net"),
        (Inches(7.6),  "Output  x' ~ p_model", ORG_D,  "New realistic sample"),
    ]
    for i, (bx, label, color, sublabel) in enumerate(flow):
        lbg = CARD_BG if color == BLUE else (GREEN_L if color == GREEN_D else ORG_L)
        add_rect(s, bx, Inches(2.2), Inches(3.0), Inches(1.55),
                 fill=lbg, line=color, lw=Pt(2))
        add_textbox(s, bx + Inches(0.15), Inches(2.3), Inches(2.7), Inches(1.3),
                    label, size=Pt(18), bold=True, color=color,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, bx, Inches(3.9), Inches(3.0), Inches(0.45),
                    sublabel, size=Pt(12), color=TEXT_S,
                    align=PP_ALIGN.CENTER, italic=True)
        if i < 2:
            add_textbox(s, bx + Inches(3.05), Inches(2.65), Inches(0.85), Inches(0.55),
                        "->", size=Pt(28), bold=True, color=TEXT_S,
                        align=PP_ALIGN.CENTER)
    add_rect(s, Inches(0.4), Inches(4.55), Inches(12.5), Inches(1.35),
             fill=GREEN_L, line=GREEN_D, lw=Pt(1.5))
    add_rect(s, Inches(0.4), Inches(4.55), Inches(0.1), Inches(1.35), fill=GREEN_D)
    add_textbox(s, Inches(0.65), Inches(4.65), Inches(12.0), Inches(0.5),
                "Key Insight: The model has not simply memorized training data.",
                size=Pt(17), bold=True, color=GREEN_D)
    add_textbox(s, Inches(0.65), Inches(5.2), Inches(12.0), Inches(0.55),
                "It has learned the underlying structure, allowing it to generate new, valid samples never seen before.",
                size=Pt(14), color=TEXT)
    return s


# ── Slide 10: Probability - Sample Space ──────────────────────────────────────
def slide_prob_sample_space(prs):
    s = new_slide(prs)
    header_bar(s, "Core Probability Theory", "Sample Space, Events and Probability")
    footer(s)
    concepts = [
        ("Sample Space   Omega",
         "The set of all possible outcomes of a random experiment.\n"
         "For a grayscale 28x28 image: Omega = [0,1]^784 -- a 784-dimensional space.",
         BLUE, CARD_BG),
        ("Event   A subset of Omega",
         "A subset of outcomes we care about.\n"
         "Example: 'the image contains a cat' or 'pixel intensity > 0.5'.",
         GREEN_D, GREEN_L),
        ("Probability Measure   P",
         "Assigns a value in [0,1] to each event:\n"
         "P(Omega) = 1,   P(empty) = 0,   0 <= P(A) <= 1 for all events A",
         ORG_D, ORG_L),
        ("Random Variable   X",
         "A function X: Omega -> R^n mapping outcomes to numeric values.\n"
         "For generative models: X represents an observation (image, text, audio).",
         PUR_D, PUR_L),
    ]
    for i, (title, body, color, light) in enumerate(concepts):
        cx = Inches(0.3) if i % 2 == 0 else Inches(6.85)
        cy = Inches(1.5) if i < 2 else Inches(4.35)
        add_rect(s, cx, cy, Inches(6.2), Inches(2.6), fill=light, line=color, lw=Pt(1.5))
        add_rect(s, cx, cy, Inches(0.1), Inches(2.6), fill=color)
        add_textbox(s, cx + Inches(0.2), cy + Inches(0.12), Inches(5.8), Inches(0.5),
                    title, size=Pt(17), bold=True, color=color)
        add_textbox(s, cx + Inches(0.2), cy + Inches(0.72), Inches(5.8), Inches(1.7),
                    body, size=Pt(13.5), color=TEXT, wrap=True)
    return s


# ── Slide 11: PDFs and PMFs ───────────────────────────────────────────────────
def slide_pdfs(prs):
    s = new_slide(prs)
    header_bar(s, "Probability Density and Mass Functions")
    footer(s)
    # PDF
    add_rect(s, Inches(0.3), Inches(1.5), Inches(6.1), Inches(5.5),
             fill=CARD_BG, line=BLUE, lw=Pt(1.2))
    add_rect(s, Inches(0.3), Inches(1.5), Inches(6.1), Inches(0.6), fill=HDR_BG)
    add_textbox(s, Inches(0.5), Inches(1.55), Inches(5.7), Inches(0.5),
                "Continuous Data -- Probability Density Function (PDF)",
                size=Pt(15), bold=True, color=WHITE)
    pdf_lines = [
        {"text": "f(x) >= 0   for all x in R^n", "size": Pt(15), "color": BLUE, "bold": True},
        {"text": "", "size": Pt(5), "color": TEXT},
        {"text": "Integral of f(x) dx = 1   (integrates to 1)", "size": Pt(15), "color": BLUE, "bold": True},
        {"text": "", "size": Pt(5), "color": TEXT},
        {"text": "P(a <= X <= b) = Integral from a to b of f(x) dx", "size": Pt(14), "color": TEXT},
        {"text": "", "size": Pt(5), "color": TEXT},
        {"text": "Gaussian (Normal) Distribution:", "size": Pt(14), "bold": True, "color": HDR_BG},
        {"text": "f(x) = (1 / sqrt(2*pi*sigma^2)) * exp( -(x-mu)^2 / (2*sigma^2) )", "size": Pt(13), "color": TEXT, "italic": True},
        {"text": "", "size": Pt(5), "color": TEXT},
        {"text": "Parameters: mu (mean), sigma^2 (variance)", "size": Pt(13), "color": TEXT_S},
        {"text": "Used as prior distribution in VAEs and diffusion models.", "size": Pt(13), "color": AMB, "italic": True},
    ]
    add_multiline(s, Inches(0.5), Inches(2.2), Inches(5.7), Inches(4.6), pdf_lines)
    # PMF
    add_rect(s, Inches(6.8), Inches(1.5), Inches(6.2), Inches(5.5),
             fill=GREEN_L, line=GREEN_D, lw=Pt(1.2))
    add_rect(s, Inches(6.8), Inches(1.5), Inches(6.2), Inches(0.6), fill=GREEN_D)
    add_textbox(s, Inches(7.0), Inches(1.55), Inches(5.8), Inches(0.5),
                "Discrete Data -- Probability Mass Function (PMF)",
                size=Pt(15), bold=True, color=WHITE)
    pmf_lines = [
        {"text": "P(X = x) >= 0   for each outcome x", "size": Pt(15), "color": GREEN_D, "bold": True},
        {"text": "", "size": Pt(5), "color": TEXT},
        {"text": "Sum over all x of P(X = x) = 1   (sums to 1)", "size": Pt(15), "color": GREEN_D, "bold": True},
        {"text": "", "size": Pt(5), "color": TEXT},
        {"text": "Categorical Distribution:", "size": Pt(14), "bold": True, "color": HDR_BG},
        {"text": "P(X = k) = p_k,   where  Sum of p_k = 1", "size": Pt(14), "color": TEXT, "italic": True},
        {"text": "", "size": Pt(5), "color": TEXT},
        {"text": "Use in text generation:", "size": Pt(14), "bold": True, "color": HDR_BG},
        {"text": "Each token is sampled from a PMF over the vocabulary.", "size": Pt(14), "color": TEXT},
        {"text": "", "size": Pt(5), "color": TEXT},
        {"text": "GPT uses a PMF at each step to choose the next token.", "size": Pt(13), "color": AMB, "italic": True},
    ]
    add_multiline(s, Inches(7.0), Inches(2.2), Inches(5.8), Inches(4.6), pmf_lines)
    return s


# ── Slide 12: Parametric Modeling ─────────────────────────────────────────────
def slide_parametric(prs):
    s = new_slide(prs)
    header_bar(s, "Parametric Modeling", "Choosing a Family of Distributions")
    footer(s)
    add_textbox(s, Inches(0.4), Inches(1.4), Inches(12.5), Inches(0.4),
                "Parametric modeling constrains p_model to a family of distributions indexed by parameters theta:",
                size=Pt(15), color=TEXT)
    eq_box(s, Inches(2.5), Inches(2.0), Inches(8.3), Inches(0.75),
           "Model family:  { p(x ; theta)  |  theta in Theta }", color=BLUE)
    families = [
        ("Gaussian",
         "theta = (mu, sigma^2)\nBell-shaped, symmetric\nUsed in VAE latent priors",
         BLUE, CARD_BG),
        ("Mixture of Gaussians",
         "theta = { (mu_k, sigma_k^2, pi_k) }\nSum of weighted Gaussians\nMore expressive than single Gaussian",
         GREEN_D, GREEN_L),
        ("Neural Network",
         "theta = weights W and biases b\nArbitrarily expressive function class\nRequires gradient-based optimization",
         PUR_D, PUR_L),
    ]
    for i, (fname, fbody, color, light) in enumerate(families):
        fx = Inches(0.3 + i * 4.35)
        add_rect(s, fx, Inches(3.1), Inches(4.1), Inches(3.75),
                 fill=light, line=color, lw=Pt(1.5))
        add_rect(s, fx, Inches(3.1), Inches(0.1), Inches(3.75), fill=color)
        add_textbox(s, fx + Inches(0.2), Inches(3.2), Inches(3.7), Inches(0.48),
                    fname, size=Pt(17), bold=True, color=color)
        add_textbox(s, fx + Inches(0.2), Inches(3.8), Inches(3.7), Inches(2.85),
                    fbody, size=Pt(13.5), color=TEXT, wrap=True)
    add_textbox(s, Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.35),
                "Tradeoff: Simple families are tractable but limited. Neural networks are expressive but harder to train.",
                size=Pt(13), color=TEXT_S, italic=True)
    return s


# ── Slide 13: MLE ─────────────────────────────────────────────────────────────
def slide_mle(prs):
    s = new_slide(prs)
    header_bar(s, "Maximum Likelihood Estimation (MLE)", "Finding the Best Parameters")
    footer(s)
    add_rect(s, Inches(0.4), Inches(1.5), Inches(12.5), Inches(2.25),
             fill=CARD_BG, line=BLUE, lw=Pt(1.5))
    add_rect(s, Inches(0.4), Inches(1.5), Inches(0.1), Inches(2.25), fill=BLUE)
    add_textbox(s, Inches(0.65), Inches(1.6), Inches(12.0), Inches(0.45),
                "Objective: Find theta* that maximizes the probability of the training data",
                size=Pt(15), color=TEXT)
    add_textbox(s, Inches(0.65), Inches(2.1), Inches(12.0), Inches(0.55),
                "theta*  =  argmax over theta  of  Product(i) p(x_i ; theta)",
                size=Pt(20), bold=True, color=BLUE, italic=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.65), Inches(2.72), Inches(12.0), Inches(0.4),
                "Equivalently:   theta*  =  argmax over theta  of  Sum(i) log p(x_i ; theta)",
                size=Pt(18), bold=True, color=HDR_BG, italic=True, align=PP_ALIGN.CENTER)
    steps = [
        ("Product -> Sum", "Products of probabilities become sums of log-probabilities:\n"
         "log(a * b) = log(a) + log(b).  Avoids numerical underflow.", BLUE, CARD_BG),
        ("Minimization", "In practice we minimize the negative log-likelihood (NLL):\n"
         "theta* = argmin  -Sum_i  log p(x_i ; theta)", GREEN_D, GREEN_L),
        ("Gaussian MLE", "For a Gaussian, MLE gives closed-form solutions:\n"
         "mu* = (1/n) Sum x_i     (sample mean)\n"
         "sigma^2* = (1/n) Sum (x_i - mu*)^2   (biased sample variance)", ORG_D, ORG_L),
        ("Deep Models", "For neural networks, use gradient descent:\n"
         "Update theta <- theta + alpha * Gradient_theta Sum log p(x_i ; theta)\n"
         "This is the foundation of ALL generative model training.", PUR_D, PUR_L),
    ]
    for i, (title, body, color, light) in enumerate(steps):
        cx = Inches(0.3) if i % 2 == 0 else Inches(6.85)
        cy = Inches(4.0) if i < 2 else Inches(5.55)
        add_rect(s, cx, cy, Inches(6.2), Inches(1.4), fill=light, line=color, lw=Pt(1))
        add_rect(s, cx, cy, Inches(0.1), Inches(1.4), fill=color)
        add_textbox(s, cx + Inches(0.2), cy + Inches(0.08), Inches(5.9), Inches(0.38),
                    title, size=Pt(14), bold=True, color=color)
        add_textbox(s, cx + Inches(0.2), cy + Inches(0.5), Inches(5.9), Inches(0.8),
                    body, size=Pt(12.5), color=TEXT, wrap=True)
    return s


# ── Slide 14: Latent Variables ────────────────────────────────────────────────
def slide_latent(prs):
    s = new_slide(prs)
    header_bar(s, "Latent Variables", "Hidden Structure in Data")
    footer(s)
    add_textbox(s, Inches(0.4), Inches(1.4), Inches(12.5), Inches(0.4),
                "Many datasets have hidden (unobserved) factors that explain variation in x.",
                size=Pt(15), color=TEXT)
    add_rect(s, Inches(0.8), Inches(1.95), Inches(11.7), Inches(1.6),
             fill=CARD_BG, line=BLUE, lw=Pt(1.5))
    add_textbox(s, Inches(1.0), Inches(2.05), Inches(11.3), Inches(0.55),
                "Generative process:   z ~ p(z)     then     x ~ p(x | z ; theta)",
                size=Pt(18), bold=True, color=BLUE, italic=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1.0), Inches(2.65), Inches(11.3), Inches(0.65),
                "Marginal:   p(x)  =  Integral of  p(x|z) * p(z)  dz     (integrate out the latent variable)",
                size=Pt(16), color=HDR_BG, italic=True, align=PP_ALIGN.CENTER)
    examples = [
        ("Face Images",
         "Latent z encodes:\n- Identity (who is the person?)\n- Expression (smile, frown)\n- Lighting, pose, age",
         BLUE, CARD_BG),
        ("Music",
         "Latent z encodes:\n- Genre (jazz, classical, pop)\n- Tempo, key, mood\n- Instrument mix",
         GREEN_D, GREEN_L),
        ("Text",
         "Latent z encodes:\n- Topic (sports, politics, science)\n- Sentiment (positive, negative)\n- Writing style, formality",
         ORG_D, ORG_L),
        ("Molecules",
         "Latent z encodes:\n- Chemical scaffold structure\n- Functional groups present\n- Drug-likeness properties",
         PUR_D, PUR_L),
    ]
    for i, (title, body, color, light) in enumerate(examples):
        ex = Inches(0.3 + i * 3.25)
        add_rect(s, ex, Inches(3.8), Inches(3.0), Inches(3.0), fill=light, line=color, lw=Pt(1.5))
        add_rect(s, ex, Inches(3.8), Inches(0.1), Inches(3.0), fill=color)
        add_textbox(s, ex + Inches(0.2), Inches(3.9), Inches(2.6), Inches(0.45),
                    title, size=Pt(14), bold=True, color=color)
        add_textbox(s, ex + Inches(0.2), Inches(4.45), Inches(2.6), Inches(2.15),
                    body, size=Pt(12.5), color=TEXT, wrap=True)
    add_textbox(s, Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.35),
                "Key insight: Latent variables are the foundation of VAEs, GANs, and Diffusion Models.",
                size=Pt(13), color=AMB, italic=True)
    return s


# ── Slide 15: Dimensionality ──────────────────────────────────────────────────
def slide_dimensionality(prs):
    s = new_slide(prs)
    header_bar(s, "The High-Dimensional Challenge",
               "Why Simple Models Fail for Images and Audio")
    footer(s)
    dims = [
        ("1D signal",       "1,000 dims",     "Audio sample (1 sec, 1kHz)",   BLUE,  CARD_BG),
        ("Grayscale image", "784 dims",        "28x28 MNIST digit",             GREEN_D, GREEN_L),
        ("Color image",     "196,608 dims",    "256x256x3 RGB photo",           ORG_D, ORG_L),
        ("HD video frame",  "6,220,800 dims",  "1920x1080x3 HD frame",          PUR_D, PUR_L),
    ]
    add_textbox(s, Inches(0.4), Inches(1.4), Inches(12.5), Inches(0.4),
                "Real-world data lives in enormous-dimensional spaces:", size=Pt(15), color=TEXT)
    for i, (label, dim_str, example, color, light) in enumerate(dims):
        bx = Inches(0.3 + i * 3.25)
        add_rect(s, bx, Inches(1.95), Inches(3.0), Inches(1.9), fill=light, line=color, lw=Pt(1.5))
        add_rect(s, bx, Inches(1.95), Inches(0.1), Inches(1.9), fill=color)
        add_textbox(s, bx + Inches(0.2), Inches(2.05), Inches(2.6), Inches(0.45),
                    label, size=Pt(14), bold=True, color=color)
        add_textbox(s, bx + Inches(0.2), Inches(2.55), Inches(2.6), Inches(0.45),
                    dim_str, size=Pt(17), bold=True, color=HDR_BG)
        add_textbox(s, bx + Inches(0.2), Inches(3.05), Inches(2.6), Inches(0.5),
                    example, size=Pt(12), color=TEXT_S, italic=True)
    add_rect(s, Inches(0.4), Inches(4.05), Inches(12.5), Inches(1.45),
             fill=PINK_L, line=RED_D, lw=Pt(1.5))
    add_rect(s, Inches(0.4), Inches(4.05), Inches(0.1), Inches(1.45), fill=RED_D)
    add_textbox(s, Inches(0.65), Inches(4.12), Inches(12.0), Inches(0.5),
                "Problem: Density estimation in high dimensions is intractable with simple models",
                size=Pt(16), bold=True, color=RED_D)
    add_textbox(s, Inches(0.65), Inches(4.68), Inches(12.0), Inches(0.7),
                "Data becomes exponentially sparse as dimensions grow. A Gaussian with 196,608 dimensions "
                "cannot capture the true structure of natural images.",
                size=Pt(14), color=TEXT, wrap=True)
    add_rect(s, Inches(0.4), Inches(5.65), Inches(12.5), Inches(1.25),
             fill=GREEN_L, line=GREEN_D, lw=Pt(1.5))
    add_rect(s, Inches(0.4), Inches(5.65), Inches(0.1), Inches(1.25), fill=GREEN_D)
    add_textbox(s, Inches(0.65), Inches(5.72), Inches(12.0), Inches(0.45),
                "Solution: Deep neural networks learn compressed latent representations",
                size=Pt(16), bold=True, color=GREEN_D)
    add_textbox(s, Inches(0.65), Inches(6.22), Inches(12.0), Inches(0.45),
                "Map high-dimensional x to low-dimensional z that captures the essential structure of the data.",
                size=Pt(14), color=TEXT)
    return s


# ── Slide 16: Taxonomy Overview ───────────────────────────────────────────────
def slide_taxonomy_overview(prs):
    s = new_slide(prs)
    header_bar(s, "Generative Model Taxonomy", "Six Major Families")
    footer(s)
    add_textbox(s, Inches(0.4), Inches(1.4), Inches(12.5), Inches(0.38),
                "Each family takes a different approach to the fundamental challenge of learning p(x):",
                size=Pt(15), color=TEXT)
    families = [
        ("Energy-Based\nModels",       "Define E(x) energy\nNo explicit normalizer",     PUR_D, PUR_L,  "~1985"),
        ("Variational\nAutoencoders",  "Encoder -> latent -> decoder\nELBO objective",    BLUE,  CARD_BG, "2013"),
        ("Generative\nAdversarial\nNetworks", "Generator vs Discriminator\nMinimax game", GREEN_D, GREEN_L, "2014"),
        ("Normalizing\nFlows",         "Invertible transforms\nExact likelihood",         ORG_D, ORG_L,  "2014"),
        ("Diffusion\nModels",          "Add then remove noise\nState of the art",         PINK_D, PINK_L, "2015"),
        ("Autoregressive\nModels",     "p(x)=Product p(x_i|x<i)\nSequential generation", TEAL,  CARD_BG, "2016"),
    ]
    for i, (name, desc, color, light, year) in enumerate(families):
        bx = Inches(0.3 + (i % 3) * 4.35)
        by = Inches(2.0) if i < 3 else Inches(4.6)
        add_rect(s, bx, by, Inches(4.1), Inches(2.35), fill=light, line=color, lw=Pt(1.5))
        add_rect(s, bx, by, Inches(4.1), Inches(0.1), fill=color)
        # Year badge
        add_rect(s, bx + Inches(3.25), by + Inches(0.12), Inches(0.7), Inches(0.4), fill=color)
        add_textbox(s, bx + Inches(3.25), by + Inches(0.14), Inches(0.7), Inches(0.36),
                    year, size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, bx + Inches(0.15), by + Inches(0.17), Inches(3.0), Inches(1.0),
                    name, size=Pt(16), bold=True, color=color)
        add_textbox(s, bx + Inches(0.15), by + Inches(1.25), Inches(3.7), Inches(0.9),
                    desc, size=Pt(12.5), color=TEXT, wrap=True)
    return s


def slide_model_family(prs, title, subtitle, approach, objective, strengths,
                       weaknesses, examples, color, light, year_str):
    s = new_slide(prs)
    header_bar(s, title, subtitle)
    footer(s)
    add_rect(s, Inches(11.85), Inches(0.12), Inches(1.3), Inches(0.68), fill=color)
    add_textbox(s, Inches(11.85), Inches(0.15), Inches(1.3), Inches(0.58),
                year_str, size=Pt(17), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Left side
    add_rect(s, Inches(0.3), Inches(1.5), Inches(7.5), Inches(2.35),
             fill=light, line=color, lw=Pt(1.5))
    add_rect(s, Inches(0.3), Inches(1.5), Inches(0.1), Inches(2.35), fill=color)
    add_textbox(s, Inches(0.5), Inches(1.58), Inches(7.1), Inches(0.42),
                "How it works", size=Pt(16), bold=True, color=color)
    add_textbox(s, Inches(0.5), Inches(2.05), Inches(7.1), Inches(1.65),
                approach, size=Pt(13.5), color=TEXT, wrap=True)
    add_rect(s, Inches(0.3), Inches(4.0), Inches(7.5), Inches(1.55),
             fill=CARD_BG, line=HDR_BG, lw=Pt(1))
    add_rect(s, Inches(0.3), Inches(4.0), Inches(0.1), Inches(1.55), fill=HDR_BG)
    add_textbox(s, Inches(0.5), Inches(4.08), Inches(7.1), Inches(0.4),
                "Training Objective", size=Pt(15), bold=True, color=HDR_BG)
    add_textbox(s, Inches(0.5), Inches(4.55), Inches(7.1), Inches(0.85),
                objective, size=Pt(13), color=TEXT, wrap=True, italic=True)
    # Right side
    add_rect(s, Inches(8.1), Inches(1.5), Inches(5.0), Inches(1.75),
             fill=GREEN_L, line=GREEN_D, lw=Pt(1))
    add_rect(s, Inches(8.1), Inches(1.5), Inches(0.08), Inches(1.75), fill=GREEN_D)
    add_textbox(s, Inches(8.28), Inches(1.55), Inches(4.6), Inches(0.4),
                "Strengths", size=Pt(14), bold=True, color=GREEN_D)
    add_textbox(s, Inches(8.28), Inches(1.98), Inches(4.6), Inches(1.15),
                strengths, size=Pt(12.5), color=TEXT, wrap=True)
    add_rect(s, Inches(8.1), Inches(3.4), Inches(5.0), Inches(1.75),
             fill=PINK_L, line=RED_D, lw=Pt(1))
    add_rect(s, Inches(8.1), Inches(3.4), Inches(0.08), Inches(1.75), fill=RED_D)
    add_textbox(s, Inches(8.28), Inches(3.45), Inches(4.6), Inches(0.4),
                "Weaknesses", size=Pt(14), bold=True, color=RED_D)
    add_textbox(s, Inches(8.28), Inches(3.88), Inches(4.6), Inches(1.15),
                weaknesses, size=Pt(12.5), color=TEXT, wrap=True)
    add_rect(s, Inches(8.1), Inches(5.3), Inches(5.0), Inches(1.75),
             fill=light, line=color, lw=Pt(1))
    add_rect(s, Inches(8.1), Inches(5.3), Inches(0.08), Inches(1.75), fill=color)
    add_textbox(s, Inches(8.28), Inches(5.35), Inches(4.6), Inches(0.4),
                "Real-World Examples", size=Pt(14), bold=True, color=color)
    add_textbox(s, Inches(8.28), Inches(5.78), Inches(4.6), Inches(1.15),
                examples, size=Pt(12.5), color=TEXT, wrap=True)
    return s


# ── Slide 23: Comparison Table ────────────────────────────────────────────────
def slide_comparison(prs):
    s = new_slide(prs)
    header_bar(s, "Generative Model Comparison", "A Quick Reference Guide")
    footer(s)
    headers = ["Model Family", "Likelihood", "Sampling", "Key Strength", "Best For"]
    col_ws = [Inches(2.3), Inches(1.8), Inches(1.7), Inches(3.3), Inches(3.4)]
    col_x = [Inches(0.15)]
    for w in col_ws[:-1]:
        col_x.append(col_x[-1] + w)
    hrow_y = Inches(1.45)
    for j, (hdr, w, x) in enumerate(zip(headers, col_ws, col_x)):
        add_rect(s, x, hrow_y, w, Inches(0.5), fill=HDR_BG)
        add_textbox(s, x + Inches(0.06), hrow_y + Inches(0.05), w - Inches(0.12), Inches(0.4),
                    hdr, size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rows = [
        ("Energy-Based",     "Unnormalized",  "MCMC (slow)",  "Flexible energy fn",   "Physics, research"),
        ("VAE",              "Lower bound",   "Fast",         "Structured latent z",   "Compression, interp."),
        ("GAN",              "Implicit",      "Fast",         "Sharp images",          "High-res synthesis"),
        ("Normalizing Flow", "Exact",         "Fast",         "Exact likelihood",      "Anomaly detection"),
        ("Diffusion",        "Lower bound",   "Slow (steps)", "Best quality",          "Image/video/audio"),
        ("Autoregressive",   "Exact",         "Sequential",   "Sequence modeling",     "Text, PixelCNN"),
    ]
    row_colors = [PUR_D, BLUE, GREEN_D, ORG_D, PINK_D, TEAL]
    lights = [PUR_L, CARD_BG, GREEN_L, ORG_L, PINK_L, CARD_BG]
    for r, (row_data, rc, rl) in enumerate(zip(rows, row_colors, lights)):
        ry = hrow_y + Inches(0.52) + r * Inches(0.82)
        bg = rl if r % 2 == 0 else RGBColor(0xF8, 0xF9, 0xFF)
        for j, (cell, w, x) in enumerate(zip(row_data, col_ws, col_x)):
            add_rect(s, x, ry, w, Inches(0.78), fill=bg,
                     line=rc if j == 0 else BORDER, lw=Pt(2 if j == 0 else 0.5))
            fc = rc if j == 0 else TEXT
            bold = (j == 0)
            add_textbox(s, x + Inches(0.06), ry + Inches(0.1), w - Inches(0.12), Inches(0.58),
                        cell, size=Pt(12.5), color=fc, align=PP_ALIGN.CENTER, bold=bold, wrap=True)
    return s


# ── Slide 24: Key Takeaways ───────────────────────────────────────────────────
def slide_takeaways(prs):
    s = new_slide(prs)
    header_bar(s, "Key Takeaways", "Chapter 1 - What Is Generative Modeling?")
    footer(s)
    takeaways = [
        (BLUE,   "1", "Generative models learn p(x) - the data distribution",
         "Unlike discriminative models p(y|x), generative models understand data structure and can create new samples."),
        (GREEN_D, "2", "MLE is the fundamental training objective",
         "We optimize theta* = argmax Sum log p(x_i ; theta). This appears in every generative model, from Gaussians to LLMs."),
        (ORG_D,  "3", "Latent variables encode hidden structure",
         "Introducing z ~ p(z) allows models to disentangle underlying factors such as style, content, and identity."),
        (PUR_D,  "4", "Six distinct families, each with unique tradeoffs",
         "VAEs, GANs, Flows, Diffusion, Autoregressive, and Energy-Based models differ in likelihood, sampling, and training."),
        (TEAL,   "5", "Deep learning makes high-dimensional generation possible",
         "Neural networks learn compressed representations that overcome the curse of dimensionality in images, audio, and text."),
    ]
    lights = [CARD_BG, GREEN_L, ORG_L, PUR_L, CARD_BG]
    for i, (color, num, title, body) in enumerate(takeaways):
        by = Inches(1.5 + i * 1.1)
        add_rect(s, Inches(0.3), by, Inches(12.7), Inches(1.0),
                 fill=lights[i], line=color, lw=Pt(0.8))
        add_rect(s, Inches(0.3), by, Inches(0.55), Inches(1.0), fill=color)
        add_textbox(s, Inches(0.3), by + Inches(0.28), Inches(0.55), Inches(0.45),
                    num, size=Pt(20), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, Inches(0.97), by + Inches(0.08), Inches(4.0), Inches(0.45),
                    title, size=Pt(15), bold=True, color=color)
        add_textbox(s, Inches(0.97), by + Inches(0.55), Inches(11.8), Inches(0.4),
                    body, size=Pt(13.5), color=TEXT, wrap=True)
    return s


# ── Slide 25: Course Roadmap ──────────────────────────────────────────────────
def slide_roadmap(prs):
    s = new_slide(prs)
    header_bar(s, "Course Roadmap", "SEAS 8525 - What is Coming Up")
    footer(s)
    weeks = [
        ("1",  "Introduction to Generative AI",      "Ch. 1",      BLUE,    True),
        ("2",  "Deep Learning Fundamentals and NLP", "Ch. 2",      BLUE,    False),
        ("3",  "Autoencoders, VAE and Latent Space", "Ch. 3",      GREEN_D, False),
        ("4",  "Transformers",                        "Ch. 9",      GREEN_D, False),
        ("5",  "Autoregressive Models",               "Ch. 5, 9",   GREEN_D, False),
        ("6",  "Multimodal LLMs, CLIP and BLIP",      "Ch. 13",     ORG_D,   False),
        ("7",  "MIDTERM EXAM",                        "-",          RED_D,   False),
        ("8",  "Diffusion Models and Video Gen",      "Ch. 8",      ORG_D,   False),
        ("9",  "GANs and Wasserstein GAN",            "Ch. 4, 10",  PINK_D,  False),
        ("10", "Normalizing Flows",                   "Ch. 6",      AMB,     False),
        ("11", "Energy-Based and Small LMs",          "Ch. 7",      PUR_D,   False),
        ("12", "Music Generation",                    "Ch. 11",     TEAL,    False),
        ("13", "World Models and RL",                 "Ch. 12, 14", TEAL,    False),
        ("14", "FINAL EXAM",                          "-",          RED_D,   False),
    ]
    for i, (num, topic, chapters, color, current) in enumerate(weeks):
        col = i // 7
        row = i % 7
        wx = Inches(0.3 + col * 6.6)
        wy = Inches(1.5 + row * 0.82)
        light = CARD_BG if not current else RGBColor(0xD8, 0xEA, 0xFF)
        lw = Pt(2.5) if current else Pt(0.8)
        add_rect(s, wx, wy, Inches(6.35), Inches(0.73), fill=light, line=color, lw=lw)
        add_rect(s, wx, wy, Inches(0.48), Inches(0.73), fill=color)
        add_textbox(s, wx, wy + Inches(0.17), Inches(0.48), Inches(0.38),
                    num, size=Pt(15), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, wx + Inches(0.55), wy + Inches(0.06), Inches(3.9), Inches(0.38),
                    topic, size=Pt(13), bold=current, color=color if current else TEXT)
        add_textbox(s, wx + Inches(0.55), wy + Inches(0.43), Inches(3.9), Inches(0.28),
                    chapters, size=Pt(11), color=TEXT_S, italic=True)
        if current:
            add_rect(s, wx + Inches(4.9), wy + Inches(0.17), Inches(1.3), Inches(0.4), fill=BLUE)
            add_textbox(s, wx + Inches(4.9), wy + Inches(0.19), Inches(1.3), Inches(0.36),
                        "You are here", size=Pt(10), bold=True, color=WHITE,
                        align=PP_ALIGN.CENTER)
    return s


# ── Build Presentation ────────────────────────────────────────────────────────
def build_presentation():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slide_title(prs)
    slide_objectives(prs)
    slide_ai_landscape(prs)
    slide_disc_vs_gen(prs)
    slide_formal_problem(prs)
    slide_applications(prs)
    slide_first_model_setup(prs)
    slide_first_model_learning(prs)
    slide_first_model_sampling(prs)
    slide_prob_sample_space(prs)
    slide_pdfs(prs)
    slide_parametric(prs)
    slide_mle(prs)
    slide_latent(prs)
    slide_dimensionality(prs)
    slide_taxonomy_overview(prs)

    # Six model family slides
    slide_model_family(prs,
        "Energy-Based Models (EBMs)", "Defining Probability via Energy",
        "Define an energy function E_theta(x): a scalar measuring how 'unlikely' a data point is.\n"
        "p_theta(x) = exp(-E_theta(x)) / Z_theta\n"
        "where Z_theta = Integral exp(-E_theta(x)) dx is the partition function.\n"
        "Lower energy corresponds to higher probability.",
        "Minimize contrastive divergence:\n"
        "L(theta) = E_data[E_theta(x)] - E_model[E_theta(x)]\n"
        "Requires MCMC sampling to estimate the model expectation.",
        "Very flexible - any neural network can be an energy function\n"
        "Unified framework encompassing many model types\n"
        "Can compose multiple energy functions (product of experts)",
        "Partition function Z_theta is intractable in high dimensions\n"
        "Slow MCMC sampling at inference time\n"
        "Training is computationally expensive",
        "Boltzmann Machines and Restricted BMs (Hinton, 1985)\n"
        "Joint Energy Models (JEM, Grathwohl et al.)\n"
        "Used in protein structure modeling",
        PUR_D, PUR_L, "~1985")

    slide_model_family(prs,
        "Variational Autoencoders (VAEs)", "Encoder-Decoder with Probabilistic Latent Space",
        "Encoder q_phi(z|x): maps input x to distribution over latent z (outputs mu and sigma).\n"
        "Decoder p_theta(x|z): reconstructs x from sampled latent z.\n"
        "Reparameterization: z = mu + sigma * epsilon,  epsilon ~ N(0,I)\n"
        "This trick allows backpropagation through the sampling step.",
        "Maximize the ELBO (Evidence Lower BOund):\n"
        "ELBO = E_q[log p_theta(x|z)] - KL(q_phi(z|x) || p(z))\n"
        "Reconstruction loss + KL regularization toward N(0,I) prior",
        "Principled probabilistic framework\n"
        "Smooth, interpretable latent space enables interpolation\n"
        "Fast sampling at inference time",
        "Blurry reconstructions (optimizes an average over samples)\n"
        "ELBO is a lower bound, not exact likelihood\n"
        "KL collapse can occur during training",
        "Image generation and style interpolation\n"
        "JTVAE for drug molecule design\n"
        "Latent space in Stable Diffusion uses a VAE encoder",
        BLUE, CARD_BG, "2013")

    slide_model_family(prs,
        "Generative Adversarial Networks (GANs)", "Learning Through Adversarial Competition",
        "Generator G(z; theta_G): maps noise z ~ p(z) to fake data.\n"
        "Discriminator D(x; theta_D): classifies real vs. fake in [0,1].\n"
        "They play a minimax game: G tries to fool D, D tries to catch G.\n"
        "At Nash equilibrium: G matches the true data distribution p_data.",
        "Minimax objective (Goodfellow et al., 2014):\n"
        "min_G max_D  E[log D(x)] + E[log(1 - D(G(z)))]\n"
        "Wasserstein GAN uses a critic + gradient penalty for stability.",
        "Very sharp, high-quality generated images\n"
        "Fast sampling at inference (single forward pass)\n"
        "No explicit likelihood required",
        "Training instability and mode collapse\n"
        "No explicit likelihood - hard to evaluate\n"
        "Sensitive to hyperparameters and architecture choices",
        "StyleGAN 2/3 (photorealistic face generation)\n"
        "CycleGAN (unpaired image-to-image translation)\n"
        "Pix2Pix (paired image translation, e.g., sketch to photo)",
        GREEN_D, GREEN_L, "2014")

    slide_model_family(prs,
        "Normalizing Flows", "Invertible Transformations for Exact Likelihood",
        "Learn an invertible mapping f_theta: z -> x where z ~ p(z) is a simple prior.\n"
        "Change-of-variables formula gives exact density:\n"
        "log p_theta(x) = log p(z) + log|det(dz/dx)|\n"
        "where z = f_theta_inverse(x).  Both forward and inverse passes are exact.",
        "Maximize exact log-likelihood:\n"
        "L(theta) = Sum_i [ log p(z_i) + log|det J_i| ]\n"
        "where J = dz/dx is the Jacobian of the inverse transform.\n"
        "Requires architectures with tractable Jacobian determinants (coupling layers).",
        "Exact likelihood computation - no approximations\n"
        "Invertible: supports both inference and generation\n"
        "Good for anomaly detection and density estimation",
        "Architectural constraints (must be invertible)\n"
        "Less expressive per parameter than unrestricted networks\n"
        "Memory-intensive at high resolutions",
        "RealNVP (NIPS 2016)\n"
        "Glow (OpenAI, high-quality face generation)\n"
        "WaveGlow (NVIDIA, high-fidelity speech synthesis)",
        ORG_D, ORG_L, "2014")

    slide_model_family(prs,
        "Diffusion Models", "Denoising as a Path to Generation",
        "Forward process q: gradually corrupt x_0 with Gaussian noise over T steps.\n"
        "x_t = sqrt(alpha_bar_t) * x_0 + sqrt(1 - alpha_bar_t) * epsilon\n"
        "After T steps, x_T is approximately pure Gaussian noise.\n"
        "Reverse process p_theta: neural network (UNet) learns to denoise step by step.",
        "Simple noise-prediction objective (DDPM, Ho et al. 2020):\n"
        "L = E[ || epsilon - epsilon_theta(x_t, t) ||^2 ]\n"
        "MSE loss only. Generation: start from x_T ~ N(0,I) and apply T denoising steps.",
        "State-of-the-art image, audio, and video quality\n"
        "Stable training with a simple MSE objective\n"
        "Highly controllable (classifier-free guidance)",
        "Slow sampling - requires many denoising steps\n"
        "Computationally expensive to train and run\n"
        "Large model sizes (UNet backbone)",
        "DALL-E 2, Imagen (text-to-image generation)\n"
        "Stable Diffusion (open-source, latent diffusion)\n"
        "Sora (video generation), AudioLDM (audio)",
        PINK_D, PINK_L, "2015")

    slide_model_family(prs,
        "Autoregressive Models", "Sequential Factorization of the Joint Distribution",
        "Factorize p(x) using the chain rule of probability:\n"
        "p(x) = p(x_1) * p(x_2|x_1) * p(x_3|x_1,x_2) * ... * p(x_n|x_1,...,x_{n-1})\n"
        "Each token/pixel is predicted conditioned on ALL previous ones.\n"
        "Transformers implement this via masked self-attention.",
        "Maximize exact log-likelihood via teacher forcing:\n"
        "L(theta) = Sum_i Sum_t log p_theta(x_it | x_i_{<t})\n"
        "Feed ground-truth previous tokens during training. Generate autoregressively at inference.",
        "Exact likelihood computation - tractable training\n"
        "Stable, scalable training via teacher forcing\n"
        "Extremely powerful for text and sequential data",
        "Sequential generation is slow for long sequences\n"
        "No compact latent space for interpolation\n"
        "Exposure bias: train on truth, generate from own outputs",
        "GPT-4, Claude, Llama (text generation)\n"
        "PixelCNN, PixelRNN (image generation)\n"
        "WaveNet (audio synthesis, Google DeepMind)",
        TEAL, CARD_BG, "2016")

    slide_comparison(prs)
    slide_takeaways(prs)
    slide_roadmap(prs)

    out = os.path.join(os.path.dirname(__file__),
                       "Week1_Chapter1_Intro_Generative_AI.pptx")
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation()
